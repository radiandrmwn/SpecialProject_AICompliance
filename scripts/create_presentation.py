"""
PPE-Watch Presentation Generator
Creates a professional PowerPoint presentation for the PPE-Watch project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "PPE-Watch"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Safety Compliance Monitoring System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(68, 68, 68)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Author info
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    author_frame = author_box.text_frame
    author_frame.text = "Radian Try Darmawan\nAsia University, Taiwan\n2nd Semester, 2024-2025"
    for para in author_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(102, 102, 102)
        para.alignment = PP_ALIGN.CENTER

def create_problem_statement_slide(prs):
    """Create problem statement slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Problem Statement"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Workplace Safety Challenges"

    # Add bullet points
    points = [
        "Construction sites have high accident rates due to PPE non-compliance",
        "Manual supervision is time-consuming, inconsistent, and prone to human error",
        "Delayed violation detection leads to preventable injuries",
        "Lack of automated tracking and accountability systems",
        "Traditional methods cannot provide real-time alerts or daily statistics"
    ]

    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

def create_objectives_slide(prs):
    """Create objectives slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Objectives"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Our Goals"

    objectives = [
        "Develop an automated PPE detection system using computer vision",
        "Implement real-time violation tracking with identity-aware counting",
        "Generate comprehensive daily compliance reports",
        "Enable supervisors to receive instant notifications via Telegram",
        "Reduce workplace accidents through proactive safety monitoring",
        "Create a scalable solution deployable across multiple sites"
    ]

    for obj in objectives:
        p = tf.add_paragraph()
        p.text = obj
        p.level = 0
        p.font.size = Pt(18)

def create_methodology_slide(prs):
    """Create methodology slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Methodology & Approach"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "YOLOv8-Based Detection Pipeline"
    tf.paragraphs[0].font.bold = True

    steps = [
        ("Person Detection", "Pre-trained YOLOv8 model on COCO dataset"),
        ("PPE Detection", "Custom YOLOv8 trained on Safety Helmet & Reflective Jacket dataset"),
        ("Multi-Object Tracking", "ByteTrack algorithm for identity consistency"),
        ("Violation Logic", "IoU-based verification: Helmet (head region) + Vest (torso region)"),
        ("Event Logging", "CSV-based storage with daily partitioning"),
        ("Report Generation", "Automated aggregation with charts and PDF summaries"),
        ("Telegram Delivery", "Real-time notifications and daily reports")
    ]

    for step_title, step_desc in steps:
        p = tf.add_paragraph()
        p.text = f"{step_title}: {step_desc}"
        p.level = 0
        p.font.size = Pt(16)

def create_key_features_slide(prs):
    """Create key features slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Features"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Technical Capabilities"
    tf.paragraphs[0].font.bold = True

    features = [
        "Dual-attribute PPE detection (helmet + reflective vest)",
        "Smart violation logic: Only flags when BOTH items missing",
        "ByteTrack integration prevents double counting",
        "Zone-based monitoring with polygon filtering",
        "Interactive Telegram bot for video uploads",
        "Automated daily reporting (CSV, PNG charts, PDF)",
        "Real-time violation screenshots for evidence",
        "Scalable architecture for multiple camera deployments"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 0
        p.font.size = Pt(17)

def create_architecture_slide(prs):
    """Create system architecture slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Architecture"

    content = slide.placeholders[1]
    tf = content.text_frame

    layers = [
        ("Input Layer", "Video uploads via Telegram, live camera feeds, video files"),
        ("Detection Layer", "YOLOv8 for person, helmet, and vest detection"),
        ("Tracking Layer", "ByteTrack for multi-object tracking with unique IDs"),
        ("Analysis Layer", "Head region detection, zone filtering, violation rules"),
        ("Storage Layer", "CSV event logging, daily partitioning"),
        ("Reporting Layer", "Statistics aggregation, chart generation (Matplotlib), PDF creation"),
        ("Delivery Layer", "Telegram Bot API for automated notifications")
    ]

    for layer_name, layer_desc in layers:
        p = tf.add_paragraph()
        p.text = f"{layer_name}: {layer_desc}"
        p.level = 0
        p.font.size = Pt(15)

def create_violation_logic_slide(prs):
    """Create violation detection logic slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Intelligent Violation Detection"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Decision Matrix"
    tf.paragraphs[0].font.bold = True

    # Add violation logic
    logic = [
        "Helmet ✓ + Vest ✓ → Compliant (Green)",
        "Helmet ✓ + Vest ✗ → Violation: 'NO VEST' (Red)",
        "Helmet ✗ + Vest ✓ → Violation: 'NO HELMET' (Red)",
        "Helmet ✗ + Vest ✗ → Violation: 'NO HELMET & NO VEST' (Red)"
    ]

    for item in logic:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)

    # Add technical details
    p = tf.add_paragraph()
    p.text = "\nTechnical Parameters:"
    p.font.bold = True
    p.font.size = Pt(18)

    params = [
        "Head region: Top 35% of person bounding box",
        "Helmet IoU threshold: 0.10 (10% overlap)",
        "Vest IoU threshold: 0.15 (15% overlap)",
        "Detection confidence: 0.25 (25% minimum)"
    ]

    for param in params:
        p = tf.add_paragraph()
        p.text = param
        p.level = 0
        p.font.size = Pt(15)

def create_results_slide(prs):
    """Create results and outcomes slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Implementation Results"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Achievements"
    tf.paragraphs[0].font.bold = True

    results = [
        "Successful real-time PPE detection with YOLOv8",
        "Accurate tracking with minimal ID switches using ByteTrack",
        "Automated daily reports delivered via Telegram",
        "Interactive bot for on-demand video processing",
        "Screenshot capture system for violation evidence",
        "Scalable architecture tested on multiple video sources",
        "Comprehensive documentation and deployment guides"
    ]

    for result in results:
        p = tf.add_paragraph()
        p.text = result
        p.level = 0
        p.font.size = Pt(18)

def create_future_work_slide(prs):
    """Create future work slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Enhancements"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Planned Improvements"
    tf.paragraphs[0].font.bold = True

    future = [
        "Deploy to actual CCTV systems with RTSP streaming",
        "Add more PPE types: Safety glasses, gloves, safety boots",
        "Implement face recognition for individual worker tracking",
        "Multi-camera fusion for comprehensive site coverage",
        "Mobile app for supervisor dashboard and alerts",
        "Integration with access control systems",
        "Cloud deployment for remote monitoring",
        "Advanced analytics: Heatmaps, trend analysis, predictive insights"
    ]

    for item in future:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)

def create_references_slide(prs):
    """Create references slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "References & Citations"

    content = slide.placeholders[1]
    tf = content.text_frame

    references = [
        "[1] Jocher, G., et al. (2023). Ultralytics YOLOv8. https://github.com/ultralytics/ultralytics",
        "[2] Zhang, Y., et al. (2022). ByteTrack: Multi-Object Tracking by Associating Every Detection Box. ECCV 2022.",
        "[3] Dataset: Safety Helmet and Reflective Jacket Dataset. https://datasetninja.com/safety-helmet-and-reflective-jacket",
        "[4] Redmon, J., & Farhadi, A. (2018). YOLOv3: An Incremental Improvement. arXiv:1804.02767.",
        "[5] Telegram Bot API Documentation. https://core.telegram.org/bots/api",
        "[6] OpenCV Library. https://opencv.org/",
        "[7] Python-Telegram-Bot. https://python-telegram-bot.org/",
        "[8] Shapely: Python package for geometric objects. https://shapely.readthedocs.io/"
    ]

    for ref in references:
        p = tf.add_paragraph()
        p.text = ref
        p.level = 0
        p.font.size = Pt(13)

def create_thank_you_slide(prs):
    """Create thank you slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Thank you text
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You!"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(54)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = RGBColor(0, 51, 102)
    thanks_para.alignment = PP_ALIGN.CENTER

    # Questions text
    questions_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    questions_frame = questions_box.text_frame
    questions_frame.text = "Questions & Discussion"
    questions_para = questions_frame.paragraphs[0]
    questions_para.font.size = Pt(32)
    questions_para.font.color.rgb = RGBColor(68, 68, 68)
    questions_para.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "PPE-Watch Project\nAsia University, Taiwan"
    for para in contact_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(102, 102, 102)
        para.alignment = PP_ALIGN.CENTER

def main():
    """Main function to create presentation"""
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create all slides
    print("Creating slides...")
    create_title_slide(prs)
    print("  [OK] Title slide")

    create_problem_statement_slide(prs)
    print("  [OK] Problem statement")

    create_objectives_slide(prs)
    print("  [OK] Objectives")

    create_methodology_slide(prs)
    print("  [OK] Methodology")

    create_key_features_slide(prs)
    print("  [OK] Key features")

    create_architecture_slide(prs)
    print("  [OK] System architecture")

    create_violation_logic_slide(prs)
    print("  [OK] Violation logic")

    create_results_slide(prs)
    print("  [OK] Results")

    create_future_work_slide(prs)
    print("  [OK] Future work")

    create_references_slide(prs)
    print("  [OK] References")

    create_thank_you_slide(prs)
    print("  [OK] Thank you slide")

    # Save presentation
    output_path = r"c:\Users\Radian Try\Documents\2nd Asia University (TW)\2nd Semester\SpecialProject\PPE_Watch_Presentation.pptx"
    prs.save(output_path)
    print(f"\n[SUCCESS] Presentation saved to: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
